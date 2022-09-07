from pptx import Presentation
import os
import win32com.client
from tkinter import *
from tkinter.filedialog import askdirectory
from tkinter.ttk import *
import os
import sys
import pandas as pd

fileIndexs={}
processFiles=[]
searchInputFiles=[]
outputFile=[]
result=[]
scale = 100
flag=0

def process(folderPath,searchInput):
    global flag
    for root, dirs, files in os.walk(folderPath):
            for name in files:
                if name.endswith('.pptx') :
                    searchInputFiles.append(os.path.join(root, name))

    for eachFile in searchInputFiles:
        count=1
        templist=[]
        prs = Presentation(eachFile) 
        for i in range(len(prs.slides)):
            for shape in prs.slides[i].shapes:
                if hasattr(shape, "text"):
                    shape.text = shape.text.lower()
                    if searchInput in shape.text:
                        templist.append(i+1)
                        fileIndexs[eachFile]=templist
                        if eachFile not in processFiles:
                            processFiles.append(eachFile)
                        name=str(eachFile).split("\\")[-1]
                        result.append((count,i+1,str(shape.text),name))
                        count=count+1
                        break
    if result != []:
        df = pd.DataFrame(result)
        df.columns = ['页码','对应页码', '内容','源文件']
        df.to_excel(folderPath+"/"+str(searchInput)+"result.xlsx", index=False)
    else:
        print("没有查到内容")
        flag=1
                
def mergePresentations(searchInputFileNames, outputFileName):

    Application = win32com.client.Dispatch("PowerPoint.Application")
    outputPresentation = Application.Presentations.Add() 
    outputPresentation.SaveAs(outputFileName)
    read_only = True
    has_title = False
    window = False

    
    for i in range(len(searchInputFileNames)):
        pb["value"] = ((i+1)/len(searchInputFileNames))*100
        currentPresentation = Application.Presentations.Open(searchInputFileNames[i],read_only,has_title,window)
        currentPresentation.Slides.Range(fileIndexs[searchInputFileNames[i]]).copy()
        Application.Presentations(outputFileName).Windows(1).Activate()    
        outputPresentation.Application.CommandBars.ExecuteMso("PasteSourceFormatting")    
        currentPresentation.Close()

    outputPresentation.save()
    outputPresentation.close()
    Application.Quit()

def savePath(path,input):
    outputFile.append(path+"/"+str(input)+"result.pptx")
        
def restart_program():
    python = sys.executable
    os.execl(python, python, * sys.argv)

def running():    
    global flag
    var = ent1.get()	    
    savePath(saveDirpath.get(),var)  
    process(workDirPath.get(),var)
    if flag==0:mergePresentations(processFiles, outputFile[0])


frameT = Tk()
frameT.geometry('500x250+400+200')
frameT.title('选择需要的PPTX文件')

frame0 = Frame(frameT)
frame0.pack(padx=10, pady=10)

frame = Frame(frameT)
frame.pack(padx=10, pady=10)  

frame_1 = Frame(frameT)
frame_1.pack(padx=10, pady=10)  

frame1 = Frame(frameT)
frame1.pack(padx=10, pady=10)

frame2 = Frame(frameT)
frame2.pack(padx=10, pady=10)

searchLabel = Label(frame0,text="搜索内容").pack(fill=X, padx=35,side=LEFT)

workDirPath = StringVar()
saveDirpath = StringVar()
searchInput=StringVar()

ent3 = Entry(frame_1, width=50, textvariable=saveDirpath).pack(fill=X, side=RIGHT)  
ent2 = Entry(frame, width=50, textvariable=workDirPath).pack(fill=X, side=RIGHT)  
ent1 = Entry(frame0, width=50, textvariable=searchInput)
ent1.pack(fill=X, side=LEFT)  

def fileopen():
    path1 = askdirectory()
    if path1:
        workDirPath.set(path1)


def fileopen_1():
    path2 = askdirectory()
    if path2:
        saveDirpath.set(path2)

btn = Button(frame, width=20, text='文件夹路径', command=fileopen).pack(fil=X, padx=10,side=LEFT)
btn_1 = Button(frame_1, width=20, text='保存路径',  command=fileopen_1).pack(fil=X, padx=10,side=LEFT)

pb = Progressbar(frameT,length=600,mode="determinate",orient=HORIZONTAL)
pb.pack(padx=10,pady=0)
pb["maximum"] = 100
pb["value"] = 0

ext = Button(frame1, width=10, text='运行', command=running).pack(fill=X, side=LEFT)
etb = Button(frame1, width=10, text='重启',  command=restart_program).pack(fill=X, padx=10)
frameT.mainloop()


